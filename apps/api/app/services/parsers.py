"""Document parsers: bytes -> page-oriented plain text.

Supported formats (v2): plain text, CSV, HTML, PDF (PyMuPDF + OCR fallback),
Word .docx, images (PNG/JPEG/TIFF via OCR), Excel/PPT (via embedded text extraction).
"""

from __future__ import annotations

import csv
import io
import logging
from dataclasses import dataclass
from pathlib import Path
from typing import TYPE_CHECKING

from bs4 import BeautifulSoup

if TYPE_CHECKING:
    import fitz

logger = logging.getLogger(__name__)


class UnsupportedDocumentError(RuntimeError):
    pass


@dataclass
class ParsedDocument:
    pages: list[str]


SUPPORTED_MIME_TYPES: dict[str, str] = {
    "text/plain": "text",
    "text/csv": "csv",
    "text/html": "html",
    "application/pdf": "pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "application/msword": "unsupported-word",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
    "application/vnd.ms-excel": "unsupported-xls",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
    "application/vnd.ms-powerpoint": "unsupported-ppt",
    "image/png": "image",
    "image/jpeg": "image",
    "image/tiff": "image",
    "image/bmp": "image",
    "image/webp": "image",
}

# Fallbacks for mislabelled uploads (Windows often reports octet-stream).
_MIME_BY_EXTENSION: dict[str, str] = {
    ".txt": "text/plain",
    ".csv": "text/csv",
    ".html": "text/html",
    ".htm": "text/html",
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".tiff": "image/tiff",
    ".tif": "image/tiff",
    ".bmp": "image/bmp",
    ".webp": "image/webp",
}


def resolve_mime_type(filename: str, declared: str) -> str:
    if declared == "application/octet-stream" or not declared:
        return _MIME_BY_EXTENSION.get(Path(filename).suffix.lower(), declared)
    return declared


def parse_file(filename: str, mime_type: str, data: bytes) -> ParsedDocument:
    mime_type = resolve_mime_type(filename, mime_type)
    kind = SUPPORTED_MIME_TYPES.get(mime_type, "unknown")
    if kind in ("unknown", "unsupported-word", "unsupported-xls", "unsupported-ppt"):
        raise UnsupportedDocumentError(
            f"Format not yet supported by the parsing pipeline: {mime_type}"
        )

    if kind == "text":
        return ParsedDocument(pages=[_decode(data)])
    if kind == "csv":
        return ParsedDocument(pages=[_parse_csv(data)])
    if kind == "html":
        return ParsedDocument(pages=[_parse_html(data)])
    if kind == "pdf":
        return ParsedDocument(pages=_parse_pdf(data))
    if kind == "docx":
        return ParsedDocument(pages=[_parse_docx(data)])
    if kind == "xlsx":
        return ParsedDocument(pages=[_parse_xlsx(data)])
    if kind == "pptx":
        return ParsedDocument(pages=[_parse_pptx(data)])
    if kind == "image":
        return ParsedDocument(pages=[_parse_image(data)])
    raise UnsupportedDocumentError(f"Unsupported MIME type: {mime_type}")


def _decode(data: bytes) -> str:
    for encoding in ("utf-8", "big5", "cp1252"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def _parse_csv(data: bytes) -> str:
    text = _decode(data)
    rows = list(csv.reader(io.StringIO(text)))
    if not rows:
        return ""
    # Render rows as tab-separated lines so table structure survives chunking.
    return "\n".join("\t".join(cell for cell in row) for row in rows)


def _parse_html(data: bytes) -> str:
    soup = BeautifulSoup(_decode(data), "html.parser")
    for tag in soup(["script", "style", "nav", "footer", "header"]):
        tag.decompose()
    text = soup.get_text("\n")
    return "\n".join(line.strip() for line in text.splitlines() if line.strip())


def _parse_pdf(data: bytes) -> list[str]:
    """Extract text with PyMuPDF; OCR pages that return empty text."""
    try:
        import fitz
    except ImportError as exc:  # pragma: no cover
        raise UnsupportedDocumentError("PyMuPDF is not installed") from exc

    doc = fitz.open(stream=data, filetype="pdf")
    pages: list[str] = []
    needs_ocr: list[int] = []

    for i, page in enumerate(doc):
        text = page.get_text().strip()
        if text:
            pages.append(text)
        else:
            pages.append("")  # placeholder
            needs_ocr.append(i)

    if needs_ocr:
        logger.info("PDF has %d/%d text-less pages — running OCR", len(needs_ocr), len(doc))
        try:
            _ocr_pdf_pages(doc, needs_ocr, pages)
        except Exception:
            logger.exception("OCR fallback failed on %d pages", len(needs_ocr))

    doc.close()
    return pages


def _ocr_pdf_pages(doc: fitz.Document, page_indices: list[int], pages_out: list[str]) -> None:
    """Render selected pages at 200 dpi, OCR them, and update pages_out."""
    from app.services.ocr import ocr_bytes

    for idx in page_indices:
        page = doc[idx]
        pix = page.get_pixmap(dpi=200)
        img_bytes = pix.tobytes("png")
        text = ocr_bytes(img_bytes)
        if text:
            pages_out[idx] = text


def _parse_image(data: bytes) -> str:
    from app.services.ocr import ocr_bytes

    text = ocr_bytes(data)
    if not text:
        logger.warning("OCR returned no text for %d-byte image", len(data))
    return text


def _parse_xlsx(data: bytes) -> str:
    """Extract text from Excel files using openpyxl."""
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts: list[str] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"[工作表: {sheet_name}]")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                parts.append("\t".join(cells))
    wb.close()
    return "\n".join(parts)


def _parse_pptx(data: bytes) -> str:
    """Extract text from PowerPoint files using python-pptx."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        slide_text: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_text.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    slide_text.append("\t".join(cells))
        if slide_text:
            parts.append(f"[投影片 {i}]")
            parts.extend(slide_text)
    return "\n".join(parts)


def _parse_docx(data: bytes) -> str:
    try:
        from docx import Document as DocxDocument
    except ImportError as exc:  # pragma: no cover
        raise UnsupportedDocumentError("python-docx is not installed") from exc

    document = DocxDocument(io.BytesIO(data))
    paragraphs = [p.text for p in document.paragraphs if p.text.strip()]
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            paragraphs.append("\t".join(cells))
    return "\n".join(paragraphs)
