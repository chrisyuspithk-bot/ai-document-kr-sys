"""Document parser unit tests."""

from __future__ import annotations

from io import BytesIO

import fitz
import pytest
from PIL import Image

from app.services.parsers import (
    ParsedDocument,
    UnsupportedDocumentError,
    parse_file,
    resolve_mime_type,
)


def test_plain_text_utf8() -> None:
    data = "仁愛堂社會服務部\n第二行".encode()
    parsed = parse_file("note.txt", "text/plain", data)
    assert parsed == ParsedDocument(pages=["仁愛堂社會服務部\n第二行"])


def test_plain_text_big5_fallback() -> None:
    data = "仁愛堂".encode("big5")
    parsed = parse_file("note.txt", "text/plain", data)
    assert parsed.pages[0] == "仁愛堂"


def test_csv_rendered_as_tsv() -> None:
    data = "姓名,部門\n陳大文,社會服務\n".encode()
    parsed = parse_file("staff.csv", "text/csv", data)
    assert "姓名\t部門" in parsed.pages[0]
    assert "陳大文\t社會服務" in parsed.pages[0]


def test_html_strips_markup() -> None:
    data = (
        "<html><head><style>p{color:red}</style></head><body>"
        "<h1>政策指引</h1><script>alert(1)</script>"
        "<p>第一段。</p><p>第二段。</p></body></html>"
    ).encode()
    parsed = parse_file("policy.html", "text/html", data)
    text = parsed.pages[0]
    assert "政策指引" in text
    assert "第一段。" in text
    assert "alert(1)" not in text
    assert "<h1>" not in text


def test_unsupported_format_raises() -> None:
    with pytest.raises(UnsupportedDocumentError):
        parse_file("deck.ppt", "application/vnd.ms-powerpoint", b"not a ppt")


def test_octet_stream_falls_back_to_extension() -> None:
    data = b"hello world"
    parsed = parse_file("note.txt", "application/octet-stream", data)
    assert parsed.pages[0] == "hello world"


def test_pdf_text_extraction() -> None:
    """PyMuPDF extracts selectable text; blank pages produce empty strings."""
    doc = fitz.open()
    page = doc.new_page(width=200, height=200)
    page.insert_text((10, 20), "Yan Oi Tong Elderly Services", fontsize=10)
    buf = doc.tobytes()
    doc.close()

    parsed = parse_file("doc.pdf", "application/pdf", buf)
    assert isinstance(parsed, ParsedDocument)
    assert len(parsed.pages) == 1
    assert "Yan Oi Tong Elderly Services" in parsed.pages[0]


def test_pdf_multi_page() -> None:
    doc = fitz.open()
    doc.new_page(width=200, height=200)
    doc.new_page(width=200, height=200)
    buf = doc.tobytes()
    doc.close()

    parsed = parse_file("doc.pdf", "application/pdf", buf)
    assert len(parsed.pages) == 2


def test_xlsx_parsing() -> None:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "長者名單"
    ws.append(["姓名", "年齡", "地區"])
    ws.append(["陳大文", 72, "屯門"])
    buf = BytesIO()
    wb.save(buf)
    wb.close()

    text = parse_file(
        "list.xlsx",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        buf.getvalue(),
    ).pages[0]
    assert "長者名單" in text
    assert "陳大文\t72\t屯門" in text


def test_pptx_parsing() -> None:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "服務計劃摘要"
    buf = BytesIO()
    prs.save(buf)

    text = parse_file(
        "deck.pptx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        buf.getvalue(),
    ).pages[0]
    assert "服務計劃摘要" in text


def test_image_ocr_route() -> None:
    """Image files route to the OCR path; verify structural output exists.

    OCR is heavy (model load) so we assert only the ParsedDocument envelope.
    """
    img = Image.new("RGB", (100, 100), (255, 255, 255))
    buf = BytesIO()
    img.save(buf, format="PNG")

    parsed = parse_file("scan.png", "image/png", buf.getvalue())
    assert isinstance(parsed, ParsedDocument)
    assert len(parsed.pages) == 1
    # OCR may return empty string for blank image — that's fine


def test_mime_resolution() -> None:
    assert resolve_mime_type("a.pdf", "application/octet-stream") == "application/pdf"
    assert resolve_mime_type("a.jpg", "application/octet-stream") == "image/jpeg"
    assert (
        resolve_mime_type("a.docx", "application/octet-stream")
        == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    )
    assert resolve_mime_type("a.txt", "") == "text/plain"
    assert resolve_mime_type("a.unknown", "") == ""
